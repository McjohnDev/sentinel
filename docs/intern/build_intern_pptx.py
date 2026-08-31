"""Generate the intern briefing PowerPoint for CBC Supervision (core agent)."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pathlib import Path

# --- CBC visual identity (DES-003 / UI brief) ---
NAVY = RGBColor(0x02, 0x06, 0x17)
NAVY_MID = RGBColor(0x0F, 0x17, 0x2A)
GOLD = RGBColor(0xD0, 0xB3, 0x35)
GOLD_DK = RGBColor(0xA8, 0x8E, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_600 = RGBColor(0x47, 0x55, 0x69)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SKY = RGBColor(0x02, 0x84, 0xC7)
EMERALD = RGBColor(0x05, 0x96, 0x69)
ROSE = RGBColor(0xE1, 0x1D, 0x48)
AMBER = RGBColor(0xD9, 0x77, 0x06)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)


def _set_run(run, text, size, color, bold=False, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _fill_line(shape, fill, line=None, line_w=Pt(1)):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_w


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill_line(s, fill, line)
    s.shadow.inherit = False
    return s


def rrect(slide, l, t, w, h, fill, line=None, adj=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    _fill_line(s, fill, line)
    try:
        s.adjustments[0] = adj
    except Exception:
        pass
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, text, size=18, color=SLATE_800, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size, color, bold, font)
    return box


def multiline(slide, l, t, w, h, lines, size=16, color=SLATE_700, bold=False, spacing=1.08):
    """lines: list of str or (str, dict)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, kw = item
        else:
            text, kw = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(kw.get("after", 6))
        p.line_spacing = spacing
        run = p.add_run()
        _set_run(
            run,
            text,
            kw.get("size", size),
            kw.get("color", color),
            kw.get("bold", bold),
            kw.get("font", "Calibri"),
        )
    return box


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def footer(slide, n, total, dark=False):
    color = SLATE_400 if dark else SLATE_500
    rect(slide, 0, Inches(7.22), W, Inches(0.28), NAVY if dark else SLATE_100)
    txt(slide, MARGIN, Inches(7.22), Inches(8), Inches(0.28),
        "CBC Supervision  ·  Intern briefing  ·  Core agent & dashboard",
        10, color if not dark else SLATE_400, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, Inches(11.4), Inches(7.22), Inches(1.4), Inches(0.28),
        f"{n}  /  {total}", 10, GOLD if dark else SLATE_600, bold=True,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def content_chrome(slide, kicker, title, n, total):
    rect(slide, 0, 0, W, H, SLATE_50)
    rect(slide, 0, 0, Inches(0.12), H, GOLD)
    rect(slide, 0, 0, W, Inches(1.18), WHITE)
    rect(slide, 0, Inches(1.16), W, Inches(0.04), GOLD)
    txt(slide, MARGIN, Inches(0.14), Inches(12), Inches(0.28),
        kicker.upper(), 11, GOLD_DK, bold=True)
    txt(slide, MARGIN, Inches(0.40), Inches(12.2), Inches(0.62),
        title, 26, NAVY, bold=True)
    footer(slide, n, total, dark=False)


def section_slide(prs, kicker, title, subtitle, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, NAVY)
    rect(slide, 0, 0, Inches(0.16), H, GOLD)
    txt(slide, Inches(0.8), Inches(2.15), Inches(11.5), Inches(0.35),
        kicker.upper(), 13, GOLD, bold=True)
    txt(slide, Inches(0.8), Inches(2.55), Inches(11.5), Inches(1.4),
        title, 40, WHITE, bold=True)
    if subtitle:
        txt(slide, Inches(0.8), Inches(4.15), Inches(11), Inches(1.2),
            subtitle, 18, SLATE_400)
    footer(slide, n, total, dark=True)
    return slide


def card(slide, l, t, w, h, title, body_lines, accent=GOLD, title_size=15):
    rrect(slide, l, t, w, h, WHITE, SLATE_200, adj=0.06)
    rect(slide, l, t, Inches(0.08), h, accent)
    txt(slide, l + Inches(0.22), t + Inches(0.14), w - Inches(0.35), Inches(0.36),
        title, title_size, NAVY, bold=True)
    multiline(
        slide,
        l + Inches(0.22),
        t + Inches(0.50),
        w - Inches(0.40),
        h - Inches(0.62),
        body_lines,
        size=13,
        color=SLATE_700,
    )


def pill(slide, l, t, w, h, label, fill, fg=WHITE):
    s = rrect(slide, l, t, w, h, fill, adj=0.5)
    txt(slide, l, t, w, h, label, 11, fg, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


def arrow_right(slide, l, t, w, h, fill=GOLD):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    _fill(s, fill)
    return s


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    TOTAL = 21

    # =====================================================================
    # 1. Title
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, GOLD)
    rrect(s, Inches(0.8), Inches(0.7), Inches(1.15), Inches(0.48), GOLD, adj=0.18)
    txt(s, Inches(0.8), Inches(0.7), Inches(1.15), Inches(0.48),
        "CBC", 16, NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.1), Inches(0.72), Inches(8), Inches(0.22),
        "Commercial Bank Cameroon", 12, SLATE_400, bold=True)
    txt(s, Inches(2.1), Inches(0.94), Inches(8), Inches(0.22),
        "IT Supervision  ·  Intern briefing", 12, GOLD)

    txt(s, Inches(0.8), Inches(2.15), Inches(11.5), Inches(1.5),
        "CBC Supervision Platform", 40, WHITE, bold=True)
    txt(s, Inches(0.8), Inches(3.55), Inches(11.5), Inches(0.7),
        "How the core agent watches a host, talks to the platform,\nand how operators see it on the dashboard.", 20, SLATE_400)

    # three chips
    for i, (lab, sub) in enumerate([
        ("01", "Core agent"),
        ("02", "Communication"),
        ("03", "Dashboard"),
    ]):
        x = Inches(0.8) + i * Inches(3.5)
        rrect(s, x, Inches(5.15), Inches(3.2), Inches(1.15), NAVY_MID, GOLD_DK, adj=0.08)
        txt(s, x + Inches(0.2), Inches(5.28), Inches(2.8), Inches(0.32), lab, 12, GOLD, bold=True)
        txt(s, x + Inches(0.2), Inches(5.58), Inches(2.8), Inches(0.5), sub, 18, WHITE, bold=True)

    footer(s, 1, TOTAL, dark=True)
    notes(s, """
Welcome the room. Introduce yourself.

Say: “This briefing explains the heart of CBC Supervision: a small program we install on each machine, how it talks securely to our central platform, and the screens operators use every morning.

We will not cover plugins, remote actions, or automation — those come later. Today is the foundation.”

Duration of the whole talk: about 25–30 minutes.
""")

    # =====================================================================
    # 2. Agenda
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "Roadmap", "What we will cover today", 2, TOTAL)
    items = [
        ("1", "Project in 3 minutes", "Why CBC is building this, and the three pieces of the system.", GOLD),
        ("2", "The core agent", "Identity, telemetry, connection. No plugins — only the built-in heartbeat.", SKY),
        ("3", "Agent ↔ platform", "Enrolment, ping, heartbeat, buffering when the network is down.", EMERALD),
        ("4", "The dashboard", "A simple tour: home, fleet, host detail, alerts.", AMBER),
    ]
    for i, (num, title, body, col) in enumerate(items):
        y = Inches(1.5) + i * Inches(1.3)
        rrect(s, MARGIN, y, Inches(12.2), Inches(1.18), WHITE, SLATE_200, adj=0.06)
        rrect(s, MARGIN + Inches(0.18), y + Inches(0.28), Inches(0.62), Inches(0.62), col, adj=0.18)
        txt(s, MARGIN + Inches(0.18), y + Inches(0.28), Inches(0.62), Inches(0.62),
            num, 18, WHITE if col != GOLD else NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.6), y + Inches(0.22), Inches(10.5), Inches(0.38), title, 20, NAVY, bold=True)
        txt(s, Inches(1.6), y + Inches(0.60), Inches(10.5), Inches(0.40), body, 15, SLATE_600)
    notes(s, """
Read the four parts. Give timings:

1. Overview — 3 min
2. Core agent — 12 min  (the main part)
3. Communication — 8 min
4. Dashboard — 5 min
Then questions.

Emphasize: “If you remember only one thing: the agent always calls the platform. The platform never opens a port on the host.”
""")

    # =====================================================================
    # 3. Why
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "1 · Overview", "Why this project exists", 3, TOTAL)
    card(s, MARGIN, Inches(1.5), Inches(5.9), Inches(5.3), "The problem", [
        "CBC runs a mixed estate: Windows, Linux, and macOS — servers and agency workstations.",
        "",
        "Supervision today is fragmented PowerShell scripts, each with its own JSON config.",
        "",
        "No single picture of health. Incidents are found late. Configuration is done machine by machine.",
        "",
        ("Banking constraint: availability, audit, and ISO 27001 / COBAC gravity.", {"bold": True, "color": SLATE_800, "size": 13}),
    ], ROSE)
    card(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(5.3), "The answer", [
        "One Python agent per host, same binary on every OS.",
        "",
        "One central platform that receives, stores, and evaluates health.",
        "",
        "One web console so operators see the fleet, not a folder of scripts.",
        "",
        "Lot 1 (this briefing) = see and collect. Lot 2 = act and automate. Lot 3 = AI assist.",
        "",
        ("Scale: ~128 hosts today, sized toward 500 for Lot 1.", {"bold": True, "color": SLATE_800, "size": 13}),
    ], EMERALD)
    notes(s, """
Keep this short. Do not dive into sprints or requirement IDs.

Say: “We are replacing a patchwork of scripts with one supervised product. The intern work sits on Lot 1: collect metrics, show them, raise alerts. We are not remotely restarting services today.”
""")

    # =====================================================================
    # 4. Three blocks
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "1 · Overview", "Three building blocks", 4, TOTAL)
    blocks = [
        ("Agent", "On every host", "A lightweight Python process. Collects CPU, RAM, disk, uptime. Identifies itself. Speaks HTTPS only.", GOLD),
        ("Platform", "In the datacentre", "FastAPI receiver. Validates messages, stores inventory and metrics, runs threshold rules, notifies by mail.", SKY),
        ("Dashboard", "In the browser", "React console for Admin, Operator, and Viewer. Situation room, fleet, host detail, alerts.", EMERALD),
    ]
    for i, (title, tag, body, col) in enumerate(blocks):
        x = MARGIN + i * Inches(4.15)
        rrect(s, x, Inches(1.55), Inches(3.95), Inches(5.15), WHITE, SLATE_200, adj=0.05)
        rect(s, x, Inches(1.55), Inches(3.95), Inches(0.12), col)
        txt(s, x + Inches(0.28), Inches(1.85), Inches(3.4), Inches(0.3), tag.upper(), 11, col, bold=True)
        txt(s, x + Inches(0.28), Inches(2.2), Inches(3.4), Inches(0.55), title, 26, NAVY, bold=True)
        txt(s, x + Inches(0.28), Inches(2.9), Inches(3.4), Inches(3.3), body, 16, SLATE_700)
    notes(s, """
Point to each block left to right: host → datacentre → operator.

Say: “Data only flows this way for Lot 1. The dashboard never talks to the agent directly. It always asks the platform.”
""")

    # =====================================================================
    # 5. Architecture
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "1 · Overview", "Architecture at a glance", 5, TOTAL)

    # flow boxes
    boxes = [
        (0.55, "Hosts", "Windows · Linux · macOS\n1 agent each"),
        (3.55, "TLS 443 →", "Outbound only\nno inbound port"),
        (6.55, "Platform", "FastAPI + Postgres\nrules + alerts"),
        (9.55, "Operators", "React dashboard\nREST + WebSocket"),
    ]
    for i, (x, title, body) in enumerate(boxes):
        fill = GOLD if i % 2 == 0 else NAVY
        fg = NAVY if i % 2 == 0 else WHITE
        rrect(s, Inches(x), Inches(1.7), Inches(2.85), Inches(2.05), fill, adj=0.08)
        txt(s, Inches(x + 0.15), Inches(1.82), Inches(2.55), Inches(0.4), title, 16, fg, bold=True, align=PP_ALIGN.CENTER)
        txt(s, Inches(x + 0.12), Inches(2.28), Inches(2.6), Inches(1.2), body, 13, fg if fill == NAVY else SLATE_800, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_right(s, Inches(x + 2.88), Inches(2.45), Inches(0.55), Inches(0.28), GOLD)

    rrect(s, MARGIN, Inches(4.1), Inches(12.2), Inches(2.7), WHITE, SLATE_200, adj=0.05)
    txt(s, Inches(0.8), Inches(4.25), Inches(11.5), Inches(0.4), "What stays out of this talk", 16, NAVY, bold=True)
    multiline(s, Inches(0.8), Inches(4.7), Inches(11.5), Inches(1.9), [
        "Plugins (CPU/disk as separate modules), log shipping, SNMP/network gear.",
        "Remote actions (restart a service, run a command) — Lot 2. The core agent rejects them.",
        "n8n automation, SSO, custom dashboards, PDF reports.",
        "Today we stay on the built-in heartbeat: identity + CPU/RAM/disk/uptime.",
    ], size=15, color=SLATE_700)
    notes(s, """
Walk left to right. Pause on “TLS 443 →”.

Say: “The agent opens the connection. Nothing listens on the host for the platform. That is a security decision, not a shortcut.”

Then read the out-of-scope box so the audience does not ask about plugins in the middle.
""")

    # =====================================================================
    # 6. Section agent
    # =====================================================================
    s = section_slide(
        prs, "Part 2", "The core agent",
        "One process per host. It knows who it is, measures the machine, and phones home.",
        6, TOTAL,
    )
    notes(s, "Short pause. “This is the longest part of the briefing. The agent is small, but every later feature sits on it.”")

    # =====================================================================
    # 7. What the agent is
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "2 · Core agent", "What the agent is — and is not", 7, TOTAL)
    rules = [
        ("Exactly one per host", "A lock file refuses a second process. Two agents on one machine would double-count and confuse alerts."),
        ("Python, self-contained", "Same logic on Windows, Linux, macOS. Packaged so the host does not need a system Python."),
        ("Low footprint", "Budget: under ~2% CPU on average and under 300 MB RAM. The agent reports its own CPU and RAM."),
        ("No inbound port", "It only makes outbound HTTPS calls to the platform (port 443). Firewalls stay closed."),
        ("Not a remote shell", "The core agent (capability L0) collects. If a remote task arrives, it is rejected."),
        ("Not a plugin host — yet", "Plugins exist in the codebase, but this briefing is the built-in collect + heartbeat path only."),
    ]
    for i, (title, body) in enumerate(rules):
        col = i % 3
        row = i // 3
        x = MARGIN + col * Inches(4.15)
        y = Inches(1.5) + row * Inches(2.65)
        card(s, x, y, Inches(3.95), Inches(2.45), title, [body], GOLD if row == 0 else SKY)
    notes(s, """
Read the six cards slowly. The intern should be able to recite the first four from memory.

Highlight: “If someone asks ‘can we SSH into the agent from the server?’ the answer is no. The agent calls us.”
""")

    # =====================================================================
    # 8. Identity
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "2 · Core agent", "How a host identifies itself", 8, TOTAL)

    identities = [
        ("machine_id", "Stable UUID stored in .machine_id. Survives reboot. This is “the same computer”."),
        ("hostname + IP", "Human name and the address it currently uses. IP can change; machine_id must not."),
        ("OS + version", "Windows / Linux / macOS and the release, so the fleet can be filtered."),
        ("agent_version", "Today 1.1.0. Lets the platform spot obsolete agents."),
        ("machine_type", "server or workstation — different offline rules (a PC after 18:00 is not an incident)."),
        ("agent_id + auth_key", "Issued by the platform at enrolment. Used on every later request."),
    ]
    for i, (k, v) in enumerate(identities):
        y = Inches(1.45) + i * Inches(0.88)
        rrect(s, MARGIN, y, Inches(12.2), Inches(0.80), WHITE, SLATE_200, adj=0.08)
        rrect(s, Inches(0.72), y + Inches(0.18), Inches(2.55), Inches(0.44), NAVY, adj=0.2)
        txt(s, Inches(0.72), y + Inches(0.18), Inches(2.55), Inches(0.44), k, 12, GOLD, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
        txt(s, Inches(3.5), y + Inches(0.16), Inches(9.0), Inches(0.50), v, 15, SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
    notes(s, """
Draw the distinction: machine_id is physical identity; agent_id is platform identity.

Analogy: machine_id is the serial number on the chassis. agent_id + auth_key is the badge the reception desk prints on day one.

If the same machine re-enrols, the platform finds it by machine_id and updates it instead of creating a duplicate.
""")

    # =====================================================================
    # 9. Telemetry families
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "2 · Core agent  ·  Telemetry", "What the core agent measures", 9, TOTAL)
    fams = [
        ("CPU", "cpu_percent\ncpu_cores", "Total utilisation over a 1-second sample, plus logical core count.", SKY),
        ("Memory", "ram_percent\nused / free / total GB", "From the OS: how full RAM is, not a guess.", GOLD),
        ("Disk", "per-mount %\nused / free / total GB", "Every filesystem of interest, not only C: or /.", EMERALD),
        ("Uptime", "uptime_seconds", "Time since last boot. Reboots become visible.", AMBER),
        ("Self", "agent_cpu_percent\nagent_ram_mb", "The agent watches itself so it cannot hide a leak.", ROSE),
        ("Context", "ip, timestamp\nconfig_version", "When, from where, and which config it thinks it has.", SLATE_700),
    ]
    for i, (title, code, body, col) in enumerate(fams):
        col_i, row = i % 3, i // 3
        x = MARGIN + col_i * Inches(4.15)
        y = Inches(1.48) + row * Inches(2.7)
        rrect(s, x, y, Inches(3.95), Inches(2.5), WHITE, SLATE_200, adj=0.06)
        rect(s, x, y, Inches(3.95), Inches(0.10), col)
        txt(s, x + Inches(0.22), y + Inches(0.22), Inches(3.5), Inches(0.36), title, 18, NAVY, bold=True)
        txt(s, x + Inches(0.22), y + Inches(0.62), Inches(3.5), Inches(0.7), code, 13, col, bold=True, font="Consolas")
        txt(s, x + Inches(0.22), y + Inches(1.4), Inches(3.5), Inches(0.9), body, 14, SLATE_700)
    notes(s, """
This is the payload of POST /api/agents/heartbeat.

Do not mention plugin metric.v1 names here. Stay on these fields.

Optional: “Services and watched files can be attached to the same heartbeat when CBC provides official lists (e.g. SWIFT). They are off by default.”
""")

    # =====================================================================
    # 10. How collected
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "2 · Core agent  ·  Telemetry", "How a sample is taken", 10, TOTAL)

    steps = [
        ("1", "Read OS", "psutil talks to the operating system: CPU, RAM, boot time, partitions."),
        ("2", "Shape the payload", "Numbers are rounded, disks listed, timestamp set in UTC, identity fields added."),
        ("3", "Send or buffer", "If the platform answers: POST heartbeat. If not: write the sample to a local disk queue."),
        ("4", "Platform decides", "Store the sample, refresh “last seen”, compare to thresholds, open or clear alerts."),
    ]
    for i, (n, title, body) in enumerate(steps):
        y = Inches(1.5) + i * Inches(1.25)
        rrect(s, MARGIN, y, Inches(12.2), Inches(1.12), WHITE, SLATE_200, adj=0.08)
        rrect(s, Inches(0.75), y + Inches(0.22), Inches(0.68), Inches(0.68), GOLD, adj=0.18)
        txt(s, Inches(0.75), y + Inches(0.22), Inches(0.68), Inches(0.68), n, 20, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.7), y + Inches(0.18), Inches(10.5), Inches(0.38), title, 18, NAVY, bold=True)
        txt(s, Inches(1.7), y + Inches(0.58), Inches(10.5), Inches(0.4), body, 15, SLATE_600)

    notes(s, """
Walk 1→4. Mention the default heartbeat interval: 30 seconds (configurable).

CPU uses a 1-second sample so the value is real, not a leftover 0%.

Self-footprint is measured across the whole interval so we do not accidentally measure the agent while it is sleeping.
""")

    # =====================================================================
    # 11. Connection
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "2 · Core agent  ·  Connection", "How the agent reaches the platform", 11, TOTAL)

    card(s, MARGIN, Inches(1.48), Inches(6.0), Inches(5.35), "Connection rules", [
        "Config key: server.url  (example https://supervision.cbc.cm)",
        "",
        "Transport: HTTPS, TLS 1.2+, port 443.",
        "",
        "tls_verify: true in production. false only in the lab with HTTP or a self-signed cert.",
        "",
        "Timeout: 10 seconds per call. Failures increment a retry counter.",
        "",
        "The agent never waits for the platform to call it. There is no agent TCP listener for this product.",
    ], GOLD)

    card(s, Inches(6.9), Inches(1.48), Inches(5.85), Inches(5.35), "What “integration” means here", [
        "Not an ESB, not a VPN plugin, not an agent management console on the host.",
        "",
        "Integration = the agent is a client of three REST endpoints:",
        "  • POST /api/agents/enroll",
        "  • POST /api/agents/ping",
        "  • POST /api/agents/heartbeat",
        "",
        "After enrolment, every call carries Authorization: <auth_key>.",
        "",
        "The dashboard is a different client (JWT users). Agents and humans do not share the same login.",
    ], SKY)
    notes(s, """
“Integration” in this briefing = how the agent is wired to the platform, not third-party tools.

If asked about Mail API / LDAP: those are platform integrations, not agent integrations. Park them.
""")

    # =====================================================================
    # 12. Enrolment
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "2 · Core agent  ·  Connection", "Enrolment — joining the fleet", 12, TOTAL)

    cols = [
        ("Admin, on the dashboard", [
            "Opens Enrol a host.",
            "Platform creates a one-time token with an expiry.",
            "Token is given to the person installing the agent (or baked into a controlled install, never into the public binary).",
        ]),
        ("Agent, on the host", [
            "Sends token + machine_id + hostname + OS + IP + machine_type.",
            "No auth_key yet — this is the only unauthenticated agent call.",
            "Retries with backoff (5s → 60s) if the platform is unreachable. The process does not exit.",
        ]),
        ("Platform, in response", [
            "Validates and consumes the token (cannot be reused).",
            "Creates or updates the agent row by machine_id.",
            "Returns agent_id + auth_key. Agent stores them in session.json and never prints the key in logs.",
        ]),
    ]
    for i, (title, lines) in enumerate(cols):
        x = MARGIN + i * Inches(4.15)
        rrect(s, x, Inches(1.48), Inches(3.95), Inches(5.35), WHITE, SLATE_200, adj=0.05)
        pill(s, x + Inches(0.22), Inches(1.68), Inches(0.55), Inches(0.32), str(i + 1), GOLD, NAVY)
        txt(s, x + Inches(0.22), Inches(2.1), Inches(3.5), Inches(0.7), title, 16, NAVY, bold=True)
        body = []
        for ln in lines:
            body.append("•  " + ln)
            body.append("")
        multiline(s, x + Inches(0.22), Inches(2.9), Inches(3.5), Inches(3.6), body, size=13, color=SLATE_700)
    notes(s, """
This is the security boundary. A random laptop cannot join: it needs a token an admin issued.

Token is single-use even if the machine was already known — that closed a reuse hole.

After this, the agent is a named member of the fleet.
""")

    # =====================================================================
    # 13. Ping vs heartbeat
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "2 · Core agent  ·  Connection", "Two rhythms: ping and heartbeat", 13, TOTAL)

    # table-like
    headers = ["", "Ping  (presence)", "Heartbeat  (telemetry)"]
    rows = [
        ["Default interval", "10 seconds", "30 seconds"],
        ["Endpoint", "POST /api/agents/ping", "POST /api/agents/heartbeat"],
        ["Body", "Empty {}", "Full metrics payload"],
        ["Purpose", "Am I still here?", "Here is my health"],
        ["Platform work", "Touch last_seen, maybe clear “offline”", "Store sample, run CPU/RAM/disk rules"],
        ["Cost", "Cheap", "Heavier (collect + write + rules)"],
    ]
    # header
    rrect(s, MARGIN, Inches(1.45), Inches(12.2), Inches(0.55), NAVY, adj=0.04)
    txt(s, Inches(0.7), Inches(1.45), Inches(3.2), Inches(0.55), headers[0], 13, GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.0), Inches(1.45), Inches(4.3), Inches(0.55), headers[1], 14, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(8.4), Inches(1.45), Inches(4.3), Inches(0.55), headers[2], 14, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    for i, row in enumerate(rows):
        y = Inches(2.05) + i * Inches(0.72)
        bg = WHITE if i % 2 == 0 else SLATE_100
        rect(s, MARGIN, y, Inches(12.2), Inches(0.72), bg)
        txt(s, Inches(0.7), y, Inches(3.2), Inches(0.72), row[0], 14, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(4.0), y, Inches(4.3), Inches(0.72), row[1], 14, SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(8.4), y, Inches(4.3), Inches(0.72), row[2], 14, SLATE_700, anchor=MSO_ANCHOR.MIDDLE)

    notes(s, """
Analogy: ping is a “present” tick on a classroom register. Heartbeat is the health form filled every 30 seconds.

Why both? Presence should be faster than metric collection. A workstation going dark is detected on the ping cadence, without forcing a full psutil sample 6 times a minute.

If ping or heartbeat returns 401/403, the agent drops its key and re-enrols.
""")

    # =====================================================================
    # 14. Degraded
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "2 · Core agent  ·  Connection", "When the platform is unreachable", 14, TOTAL)

    # timeline
    phases = [
        ("Online", "Collect → send\nPlatform ACKs\nBuffer empty", EMERALD),
        ("Outage", "Collect continues\nSamples written to\nqueue.jsonl on disk", AMBER),
        ("Caps", "Keep at most\n500 MB or 24 hours\nOldest records drop", ROSE),
        ("Recovery", "Checkout the queue\nReplay heartbeats\nThen resume live send", SKY),
    ]
    for i, (title, body, col) in enumerate(phases):
        x = MARGIN + i * Inches(3.15)
        rrect(s, x, Inches(1.55), Inches(2.95), Inches(3.15), WHITE, SLATE_200, adj=0.08)
        rect(s, x, Inches(1.55), Inches(2.95), Inches(0.1), col)
        txt(s, x + Inches(0.18), Inches(1.8), Inches(2.6), Inches(0.45), title, 18, NAVY, bold=True)
        txt(s, x + Inches(0.18), Inches(2.35), Inches(2.6), Inches(2.1), body, 14, SLATE_700)
        if i < 3:
            arrow_right(s, x + Inches(2.96), Inches(2.9), Inches(0.18), Inches(0.18), GOLD)

    rrect(s, MARGIN, Inches(4.95), Inches(12.2), Inches(1.85), NAVY, adj=0.06)
    txt(s, Inches(0.8), Inches(5.15), Inches(11.5), Inches(0.4), "Why this matters in a bank", 16, GOLD, bold=True)
    txt(s, Inches(0.8), Inches(5.55), Inches(11.5), Inches(1.0),
        "A platform restart or WAN blip must not erase the last hours of CPU/disk history. "
        "The agent stays up, keeps sampling, and delivers the gap when the path is back. "
        "Replay is crash-safe: a kill during flush does not delete the batch.",
        15, WHITE)
    notes(s, """
Name: “degraded mode” / store-and-forward.

Enrolment failures also buffer: the agent does not sit idle; it still collects locally.

Do not over-promise: if the cap is hit, oldest data is dropped. 24 h / 500 MB is the contract.
""")

    # =====================================================================
    # 15. Section communication
    # =====================================================================
    s = section_slide(
        prs, "Part 3", "Communication with the platform",
        "A small, explicit conversation. The agent speaks. The platform answers.",
        15, TOTAL,
    )
    notes(s, "Transition: “We now look at the same facts from the wire: URLs, payloads, and what the server does.”")

    # =====================================================================
    # 16. Communication map
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "3 · Communication", "The conversation — four calls", 16, TOTAL)

    calls = [
        ("POST /api/agents/enroll", "Once (or after auth loss)", "Token + identity", "agent_id, auth_key", "Unauthenticated"),
        ("POST /api/agents/ping", "Every ~10 s", "Empty body", "status, server_time", "auth_key"),
        ("POST /api/agents/heartbeat", "Every ~30 s", "CPU/RAM/disk/uptime", "optional config push", "auth_key"),
        ("POST /api/agents/config/ack", "When config version changes", "version number", "200 OK", "auth_key"),
    ]
    rrect(s, MARGIN, Inches(1.42), Inches(12.2), Inches(0.48), NAVY, adj=0.04)
    for i, h in enumerate(["Call", "When", "Agent sends", "Platform returns", "Auth"]):
        txt(s, Inches(0.65 + i * 2.45), Inches(1.42), Inches(2.35), Inches(0.48), h, 12, GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    for r, row in enumerate(calls):
        y = Inches(1.95) + r * Inches(0.95)
        bg = WHITE if r % 2 == 0 else SLATE_100
        rect(s, MARGIN, y, Inches(12.2), Inches(0.95), bg)
        for c, val in enumerate(row):
            font = "Consolas" if c == 0 else "Calibri"
            size = 11 if c == 0 else 13
            colr = NAVY if c == 0 else SLATE_700
            txt(s, Inches(0.65 + c * 2.45), y, Inches(2.4), Inches(0.95), val, size, colr, bold=(c == 0),
                anchor=MSO_ANCHOR.MIDDLE, font=font)

    txt(s, MARGIN, Inches(5.9), Inches(12.2), Inches(0.9),
        "Dashboard traffic is separate: operators use JWT on /api/agents, /api/alerts, WebSocket. "
        "The agent never uses a user password.",
        15, SLATE_600)
    notes(s, """
Skip plugin ingest (/api/ingest/metrics, /api/ingest/logs) unless someone asks — then say “exists, out of scope today.”

Config ack: the heartbeat response can carry a new YAML overlay. The agent merges it, then acks the version. Core story: central config, not SSH to edit files.
""")

    # =====================================================================
    # 17. Sequences
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "3 · Communication", "Two sequences to remember", 17, TOTAL)

    card(s, MARGIN, Inches(1.48), Inches(6.05), Inches(5.35), "A. First contact (enrolment)", [
        "1. Admin issues a one-time token.",
        "2. Agent POST /enroll with token + machine_id.",
        "3. Platform checks token, upserts the host.",
        "4. Platform returns agent_id + auth_key.",
        "5. Agent writes session.json and starts ping/heartbeat.",
        "",
        "If the token is wrong or expired → 4xx, agent retries later. No silent join.",
    ], GOLD)

    card(s, Inches(6.9), Inches(1.48), Inches(5.85), Inches(5.35), "B. Steady state (every 30 s)", [
        "1. Agent samples CPU / RAM / disk.",
        "2. POST /heartbeat with Authorization.",
        "3. Platform stores the row, updates last_communication.",
        "4. Rules compare values to warning/critical thresholds.",
        "5. Alerts open, stay, or clear. Mail if configured.",
        "6. Dashboard sees the new last-seen and gauges (REST / live WS).",
        "",
        "Between heartbeats, ping keeps last_seen fresh.",
    ], SKY)
    notes(s, """
You can mime this with your hands: left hand = agent, right hand = platform.

Sequence A happens rarely. Sequence B is the product.

Mention availability windows only if asked: workstations have office-hour calendars so overnight silence is not a critical alert.
""")

    # =====================================================================
    # 18. Heartbeat journey
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "3 · Communication", "Inside the platform, after a heartbeat", 18, TOTAL)

    steps = [
        ("Receive", "FastAPI authenticates the auth_key and parses the JSON body."),
        ("Inventory", "Agent row: last_communication, IP, self CPU/RAM, config version."),
        ("Store", "Heartbeat row in PostgreSQL. Samples also written to the time-series store."),
        ("Rules", "CPU / RAM / disk vs global or per-host thresholds. Offline if last_seen too old."),
        ("Notify", "New critical/major alerts can go out on the CBC Mail Service API."),
        ("Show", "Dashboard KPIs, fleet status, host charts, alert list refresh."),
    ]
    for i, (title, body) in enumerate(steps):
        col, row = i % 3, i // 3
        x = MARGIN + col * Inches(4.15)
        y = Inches(1.5) + row * Inches(2.6)
        rrect(s, x, y, Inches(3.95), Inches(2.4), WHITE, SLATE_200, adj=0.06)
        rrect(s, x + Inches(0.22), y + Inches(0.22), Inches(0.48), Inches(0.48), GOLD, adj=0.2)
        txt(s, x + Inches(0.22), y + Inches(0.22), Inches(0.48), Inches(0.48), str(i + 1), 16, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.85), y + Inches(0.28), Inches(2.9), Inches(0.4), title, 18, NAVY, bold=True)
        txt(s, x + Inches(0.22), y + Inches(0.9), Inches(3.5), Inches(1.25), body, 14, SLATE_700)
    notes(s, """
This is the only “backend internals” slide. Keep it operational, not framework trivia.

Postgres = source of truth for “who is the host” and “what is open”. Time-series = charts over hours/days.
""")

    # =====================================================================
    # 19. Section dashboard
    # =====================================================================
    s = section_slide(
        prs, "Part 4", "The central dashboard",
        "A situation room for operators — not a YAML editor.",
        19, TOTAL,
    )
    notes(s, "“Last part, five minutes. Show the four screens an intern must recognise.”")

    # =====================================================================
    # 20. Dashboard simple
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "4 · Dashboard", "Four screens, one job", 20, TOTAL)

    screens = [
        ("Home", "Tableau de bord", "Parc online/offline, open criticals, mail channel health, load trend, triage list. Admin can enrol a host from here."),
        ("Fleet", "Parc / Agents", "Every enrolled host: OS, status, last seen, CPU/RAM/disk gauges. Filter offline. Open a host."),
        ("Host", "Détail agent", "Identity, live gauges, metric history, this host’s alerts, thresholds. This is “why is this machine red?”"),
        ("Alerts", "Alertes", "Open / acknowledged / resolved. Severity. Ack with comment. The operator’s inbox."),
    ]
    for i, (en, fr, body) in enumerate(screens):
        y = Inches(1.45) + i * Inches(1.3)
        rrect(s, MARGIN, y, Inches(12.2), Inches(1.18), WHITE, SLATE_200, adj=0.08)
        rrect(s, Inches(0.75), y + Inches(0.22), Inches(2.35), Inches(0.74), NAVY, adj=0.12)
        txt(s, Inches(0.75), y + Inches(0.22), Inches(2.35), Inches(0.42), en, 16, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(0.75), y + Inches(0.58), Inches(2.35), Inches(0.32), fr, 11, GOLD, align=PP_ALIGN.CENTER)
        txt(s, Inches(3.35), y + Inches(0.22), Inches(9.1), Inches(0.78), body, 15, SLATE_700, anchor=MSO_ANCHOR.MIDDLE)

    notes(s, """
If you have a live demo, click these four routes. If not, stay on this slide.

Sidebar groups: Exploiter (these four) is Lot 1 daily work. Analyser / Automatiser / Configurer are later or admin.

Roles:
• Admin — enrol, users, settings, revoke
• Operator — triage, ack, diagnose
• Viewer — read only

Login is a CBC user (JWT), not an agent key.
""")

    # =====================================================================
    # 21. Close
    # =====================================================================
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, GOLD)
    txt(s, Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.3), "TO TAKE AWAY".upper(), 12, GOLD, bold=True)
    txt(s, Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7), "Four sentences", 32, WHITE, bold=True)

    takeaways = [
        "CBC Supervision replaces scattered scripts with one agent, one platform, one console.",
        "The core agent identifies the host and measures CPU, RAM, disk, uptime — then phones home on TLS 443.",
        "Enrolment issues a badge (agent_id + auth_key). Ping says “present”. Heartbeat says “here is my health”. If the path is down, samples wait on disk.",
        "The dashboard is how humans see that health. It never talks to the agent; it always asks the platform.",
    ]
    for i, line in enumerate(takeaways):
        y = Inches(1.7) + i * Inches(1.05)
        rrect(s, Inches(0.8), y, Inches(0.55), Inches(0.55), GOLD, adj=0.2)
        txt(s, Inches(0.8), y, Inches(0.55), Inches(0.55), str(i + 1), 18, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.55), y, Inches(10.8), Inches(0.9), line, 16, WHITE, anchor=MSO_ANCHOR.MIDDLE)

    txt(s, Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.7),
        "Questions?  ·  Out of scope today: plugins, logs, remote actions, n8n.",
        16, SLATE_400)
    footer(s, 21, TOTAL, dark=True)
    notes(s, """
Read the four sentences slowly. Stop. Ask for questions.

If a question is about plugins: “Next chapter — collectors become modules with a manifest. The heartbeat path you saw stays the backbone.”

If a question is about actions: “Lot 2. L0 agents reject task.v1 on purpose.”

Thank the room.
""")

    out = Path(__file__).resolve().parent / "CBC_Supervision_Intern_Briefing.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build()
    print(path)
