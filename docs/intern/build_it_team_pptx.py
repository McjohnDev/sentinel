"""Présentation IT — CBC Supervision, en français, à deux voix."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pathlib import Path

NAVY = RGBColor(0x02, 0x06, 0x17)
NAVY_MID = RGBColor(0x0F, 0x17, 0x2A)
GOLD = RGBColor(0xD0, 0xB3, 0x35)
GOLD_DK = RGBColor(0xA8, 0x8E, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_50 = RGBColor(0xF8, 0xFA, 0xFC)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_600 = RGBColor(0x47, 0x55, 0x69)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SKY = RGBColor(0x02, 0x84, 0xC7)
EMERALD = RGBColor(0x05, 0x96, 0x69)
ROSE = RGBColor(0xE1, 0x1D, 0x48)
AMBER = RGBColor(0xD9, 0x77, 0x06)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
TOTAL = 18


def _set_run(run, text, size, color, bold=False, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font


def _fill_line(shape, fill, line=None, line_w=Pt(1)):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_w


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill_line(s, fill, line)
    s.shadow.inherit = False
    return s


def rrect(slide, l, t, w, h, fill, line=None, adj=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    _fill_line(s, fill, line)
    try:
        s.adjustments[0] = adj
    except Exception:
        pass
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, text, size=18, color=SLATE_800, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size, color, bold, font)
    return box


def multiline(slide, l, t, w, h, lines, size=16, color=SLATE_700, bold=False, spacing=1.08):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, kw = item
        else:
            text, kw = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(kw.get("after", 6))
        p.line_spacing = spacing
        run = p.add_run()
        _set_run(run, text, kw.get("size", size), kw.get("color", color), kw.get("bold", bold), kw.get("font", "Calibri"))
    return box


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def footer(slide, n, dark=False):
    color = SLATE_400 if dark else SLATE_500
    rect(slide, 0, Inches(7.22), W, Inches(0.28), NAVY if dark else SLATE_100)
    txt(slide, MARGIN, Inches(7.22), Inches(9.5), Inches(0.28),
        "CBC Supervision  ·  Restitution IT  ·  Ce qui est livré",
        10, color if not dark else SLATE_400, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, Inches(11.4), Inches(7.22), Inches(1.4), Inches(0.28),
        f"{n}  /  {TOTAL}", 10, GOLD if dark else SLATE_600, bold=True,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def speaker_pill(slide, label, fill=GOLD, fg=NAVY):
    rrect(slide, Inches(11.15), Inches(0.18), Inches(1.65), Inches(0.32), fill, adj=0.5)
    txt(slide, Inches(11.15), Inches(0.18), Inches(1.65), Inches(0.32),
        label, 10, fg, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def content_chrome(slide, kicker, title, n, who="VOUS"):
    rect(slide, 0, 0, W, H, SLATE_50)
    rect(slide, 0, 0, Inches(0.12), H, GOLD)
    rect(slide, 0, 0, W, Inches(1.18), WHITE)
    rect(slide, 0, Inches(1.16), W, Inches(0.04), GOLD)
    txt(slide, MARGIN, Inches(0.14), Inches(10), Inches(0.28), kicker.upper(), 11, GOLD_DK, bold=True)
    txt(slide, MARGIN, Inches(0.40), Inches(10.4), Inches(0.62), title, 26, NAVY, bold=True)
    speaker_pill(slide, who)
    footer(slide, n, dark=False)


def section_slide(prs, kicker, title, subtitle, n, who):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, NAVY)
    rect(slide, 0, 0, Inches(0.16), H, GOLD)
    txt(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.35), kicker.upper(), 13, GOLD, bold=True)
    txt(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.5), title, 36, WHITE, bold=True)
    if subtitle:
        txt(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(1.4), subtitle, 18, SLATE_400)
    rrect(slide, Inches(0.8), Inches(5.7), Inches(2.2), Inches(0.42), GOLD, adj=0.4)
    txt(slide, Inches(0.8), Inches(5.7), Inches(2.2), Inches(0.42), who, 12, NAVY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, n, dark=True)
    return slide


def card(slide, l, t, w, h, title, body_lines, accent=GOLD, title_size=15):
    rrect(slide, l, t, w, h, WHITE, SLATE_200, adj=0.06)
    rect(slide, l, t, Inches(0.08), h, accent)
    txt(slide, l + Inches(0.22), t + Inches(0.14), w - Inches(0.35), Inches(0.36), title, title_size, NAVY, bold=True)
    multiline(slide, l + Inches(0.22), t + Inches(0.50), w - Inches(0.40), h - Inches(0.62),
              body_lines, size=13, color=SLATE_700)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1. Titre
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, GOLD)
    rrect(s, Inches(0.8), Inches(0.65), Inches(1.15), Inches(0.44), GOLD, adj=0.18)
    txt(s, Inches(0.8), Inches(0.65), Inches(1.15), Inches(0.44), "CBC", 16, NAVY, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.1), Inches(0.68), Inches(9), Inches(0.20), "Commercial Bank Cameroun", 12, SLATE_400, bold=True)
    txt(s, Inches(2.1), Inches(0.90), Inches(9), Inches(0.20), "Direction des systèmes d'information", 12, GOLD)
    txt(s, Inches(0.8), Inches(2.05), Inches(11.7), Inches(1.35),
        "Le parc n'est plus un dossier\nde scripts.", 36, WHITE, bold=True)
    txt(s, Inches(0.8), Inches(3.55), Inches(11.5), Inches(0.9),
        "CBC Supervision — ce qui est réellement livré.\nUne restitution à deux voix : l'agent sur l'hôte, la plateforme au centre.",
        18, SLATE_400)
    rrect(s, Inches(0.8), Inches(5.15), Inches(5.5), Inches(1.25), NAVY_MID, GOLD_DK, adj=0.08)
    txt(s, Inches(1.0), Inches(5.28), Inches(5.1), Inches(0.28), "VOUS", 11, GOLD, bold=True)
    txt(s, Inches(1.0), Inches(5.55), Inches(5.1), Inches(0.6), "Application centrale\nréception, règles, console, alertes", 15, WHITE)
    rrect(s, Inches(6.6), Inches(5.15), Inches(5.7), Inches(1.25), NAVY_MID, GOLD_DK, adj=0.08)
    txt(s, Inches(6.8), Inches(5.28), Inches(5.3), Inches(0.28), "INTERNE", 11, GOLD, bold=True)
    txt(s, Inches(6.8), Inches(5.55), Inches(5.3), Inches(0.6), "Agent sur l'hôte\nenrôlement, battement, métriques", 15, WHITE)
    footer(s, 1, dark=True)
    notes(s, """
VOUS — ouverture (45 s).

Debout. Regard salle. Ne lisez pas le titre.

« Ce matin, si le serveur métier d'une agence s'arrête à 7 h 40, comment le savons-nous ?
Parce qu'un utilisateur appelle — ou parce que la machine nous a déjà parlé ? »

Présentez l'interne par son prénom. « Elle va vous montrer ce qui tourne sur le poste.
Moi, ce qui se passe une fois le message arrivé chez nous. On se passe le relais. »

Durée promise : 25 minutes + questions. Tenez-la.
""")

    # 2. Pacte
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "Le pacte", "Comment on va parler — et ce qu'on ne vend pas", 2, "ENSEMBLE")
    items = [
        ("25 minutes", "Pas un catalogue de sprints. Une histoire : hôte → fil → écran."),
        ("Deux voix", "L'interne tient l'agent. Vous tenez le centre. On se coupe net, on se relance."),
        ("Honnêteté", "Dix points livrés, démontrables. Pas d'écran qui affirme ce qui n'existe pas."),
        ("Une question pour vous", "À la fin : ce dont l'IT a besoin de nous, et ce dont nous avons besoin de vous."),
    ]
    for i, (title, body) in enumerate(items):
        y = Inches(1.48) + i * Inches(1.28)
        rrect(s, MARGIN, y, Inches(12.2), Inches(1.15), WHITE, SLATE_200, adj=0.06)
        rrect(s, Inches(0.75), y + Inches(0.28), Inches(0.55), Inches(0.55), GOLD, adj=0.18)
        txt(s, Inches(0.75), y + Inches(0.28), Inches(0.55), Inches(0.55), str(i + 1), 16, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.55), y + Inches(0.22), Inches(10.8), Inches(0.38), title, 18, NAVY, bold=True)
        txt(s, Inches(1.55), y + Inches(0.62), Inches(10.8), Inches(0.38), body, 15, SLATE_600)
    notes(s, """
ENSEMBLE — 40 s. VOUS lisez les quatre lignes. INTERNE hoche la tête sur « deux voix ».

Phrase de clôture : « Si une question porte sur l'agent, elle lui revient. Si elle porte sur la console, elle me revient. »
""")

    # 3. Problème
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "Pourquoi on construit ça", "Superviser un parc bancaire dans le noir", 3, "VOUS")
    card(s, MARGIN, Inches(1.48), Inches(6.05), Inches(5.35), "Ce que l'IT vivait", [
        "Un parc hétérogène : Windows, Linux, parfois macOS — serveurs et postes d'agence.",
        "",
        "Des scripts PowerShell dispersés, chacun avec son JSON, sa cadence, son oubli.",
        "",
        "Pas de photo unique. L'incident arrive par téléphone. La config se fait machine par machine.",
        "",
        ("En banque, ça n'est pas un inconfort : c'est de la disponibilité, de l'audit, de la gravité COBAC.",
         {"bold": True, "color": SLATE_800, "size": 13}),
    ], ROSE)
    card(s, Inches(6.9), Inches(1.48), Inches(5.85), Inches(5.35), "Ce que ça coûte vraiment", [
        "MTTR trop long : on cherche encore l'hôte pendant que le métier attend.",
        "",
        "Faux calme la nuit : un poste éteint n'est pas un serveur mort — encore faut-il le distinguer.",
        "",
        "Dette opérationnelle : personne n'ose toucher le script de l'autre.",
        "",
        "Et surtout : pas de preuve. Pas d'historique. Pas de « depuis quand ». ",
    ], AMBER)
    notes(s, """
VOUS — 2 min. Ton grave, pas agressif envers l'existant.

« Personne ici n'a mal fait son travail. L'outil n'était plus à la taille du parc. »

Ne citez pas FS0. Ne citez pas GitHub. Restez métier.
""")

    # 4. Promesse
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "La promesse", "Un agent. Un centre. Une console.", 4, "VOUS")
    blocks = [
        ("Sur l'hôte", "L'agent", "Un processus léger. Il sait qui il est. Il mesure. Il appelle. Il n'écoute personne.", GOLD),
        ("Au centre", "La plateforme", "Elle authentifie, stocke, compare aux seuils, notifie, journalise. Source de vérité.", SKY),
        ("Devant l'opérateur", "La console", "Tableau de bord, parc, alerte, paramétrage. Jamais un accès direct à la machine.", EMERALD),
    ]
    for i, (tag, title, body, col) in enumerate(blocks):
        x = MARGIN + i * Inches(4.15)
        rrect(s, x, Inches(1.55), Inches(3.95), Inches(4.35), WHITE, SLATE_200, adj=0.05)
        rect(s, x, Inches(1.55), Inches(3.95), Inches(0.10), col)
        txt(s, x + Inches(0.24), Inches(1.82), Inches(3.45), Inches(0.28), tag.upper(), 11, col, bold=True)
        txt(s, x + Inches(0.24), Inches(2.18), Inches(3.45), Inches(0.55), title, 24, NAVY, bold=True)
        txt(s, x + Inches(0.24), Inches(2.85), Inches(3.45), Inches(2.6), body, 16, SLATE_700)
    txt(s, MARGIN, Inches(6.1), Inches(12.2), Inches(0.8),
        "Règle d'or : l'agent ouvre la connexion (HTTPS). Rien n'écoute sur le poste pour nous. Le firewall de l'hôte reste fermé.",
        15, SLATE_700)
    notes(s, """
VOUS — 90 s. Pointez gauche → droite.

« Les données ne vont que dans ce sens pour l'instant. La console ne parle jamais à l'agent. Elle nous demande, à nous. »

Puis : « Je te passe la parole. Tu es sur la machine. »
""")

    # 5. Section agent
    s = section_slide(
        prs, "Relais 1", "L'agent — ce qui vit sur la machine",
        "Elle va vous raconter comment un poste rejoint le parc, reste vivant, et se tait proprement.",
        5, "INTERNE",
    )
    notes(s, """
INTERNE — 10 s. Sourire. « Merci. On quitte le datacentre. On est sur le poste. »

Ne commencez pas par l'architecture. Commencez par : « Un seul programme. Un seul processus. »
""")

    # 6. Ce qu'est l'agent
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "L'agent", "Ce qu'il est — et ce qu'il refuse d'être", 6, "INTERNE")
    rules = [
        ("Un seul par hôte", "Un verrou empêche un second processus. Deux agents, c'est compter deux fois et alerter faux."),
        ("Identité stable", "Un machine_id survit au redémarrage. Réinstaller met à jour la fiche — ça ne crée pas un jumeau."),
        ("Il appelle, il n'écoute pas", "HTTPS sortant. Pas de port ouvert. La plateforme ne « ping » pas l'hôte."),
        ("Pas un shell distant", "Collecter, battre, se désinstaller. Pas d'ordre de redémarrer un service depuis ici."),
    ]
    for i, (title, body) in enumerate(rules):
        col, row = i % 2, i // 2
        x = MARGIN + col * Inches(6.25)
        y = Inches(1.48) + row * Inches(2.55)
        card(s, x, y, Inches(6.05), Inches(2.38), title, [body], GOLD if row == 0 else SKY)
    notes(s, """
INTERNE — 2 min 30. Lent. Les quatre cartes sont le contrat.

Si on vous coupe : « On y revient — d'abord le badge. »

Analogie : machine_id = numéro de châssis. agent_id + clé = badge imprimé à l'accueil le jour 1.
""")

    # 7. Enrôlement
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "L'agent  ·  Point 1", "Enrôlement : le poste obtient son badge", 7, "INTERNE")
    cols = [
        ("L'admin, au centre", [
            "Émet un jeton depuis Paramètres.",
            "Usage unique, avec échéance.",
            "On ne le grave jamais dans le binaire public.",
        ]),
        ("L'agent, sur l'hôte", [
            "Envoie jeton + machine_id + nom + OS + IP.",
            "C'est le seul appel sans badge.",
            "Un refus n'écrit rien sur le disque.",
        ]),
        ("La plateforme répond", [
            "Consomme le jeton — il ne se réutilise pas.",
            "Crée ou met à jour la fiche (même machine_id).",
            "Rend agent_id + clé. La clé ne s'imprime pas dans les journaux.",
        ]),
    ]
    for i, (title, lines) in enumerate(cols):
        x = MARGIN + i * Inches(4.15)
        rrect(s, x, Inches(1.48), Inches(3.95), Inches(5.35), WHITE, SLATE_200, adj=0.05)
        rrect(s, x + Inches(0.22), Inches(1.68), Inches(0.48), Inches(0.32), GOLD, adj=0.4)
        txt(s, x + Inches(0.22), Inches(1.68), Inches(0.48), Inches(0.32), str(i + 1), 12, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.22), Inches(2.15), Inches(3.5), Inches(0.7), title, 16, NAVY, bold=True)
        body = []
        for ln in lines:
            body.append("•  " + ln)
            body.append("")
        multiline(s, x + Inches(0.22), Inches(2.95), Inches(3.5), Inches(3.5), body, size=14, color=SLATE_700)
    notes(s, """
INTERNE — 2 min.

« Un laptop trouvé dans un tiroir ne rejoint pas le parc. Il lui faut un jeton qu'un admin a émis. »

Démo possible : python agent/src/cli.py enroll — ou le script de démo. Si la démo plante, cette slide suffit.
""")

    # 8. Battement
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "L'agent  ·  Points 5 et 7", "Il frappe. S'il se tait, on le sait.", 8, "INTERNE")
    card(s, MARGIN, Inches(1.48), Inches(6.05), Inches(5.35), "Le battement", [
        "Toutes les ~30 secondes : CPU, RAM, disques, faits d'hôte relus.",
        "",
        "La première mesure CPU n'est plus un faux 100 % au démarrage — piège classique de psutil.",
        "",
        "L'adresse IP repart à chaque cycle : un poste DHCP ne reste pas figé à l'IP du lancement.",
        "",
        "La plateforme répond dans le battement : qui elle croit qu'il est, l'écart d'horloge, depuis quand le silence a cessé.",
    ], GOLD)
    card(s, Inches(6.9), Inches(1.48), Inches(5.85), Inches(5.35), "La coupure", [
        "Si le centre est injoignable : recul progressif (5 s, 10 s… plafond). Tout le parc ne martèle pas à la seconde.",
        "",
        "Seul un 401 = perte d'identité. Un 403 (révoqué) ne se réenrôle pas tout seul — ce serait rentrer par la fenêtre.",
        "",
        "Honnêteté : on reprend le contact, on ne rejoue pas encore 24 h d'historique. Ça viendra, proprement.",
    ], SKY)
    notes(s, """
INTERNE — 2 min 30. C'est votre slide la plus technique. Une analogie :

« Le battement, c'est la fiche santé. Pas un ping ICMP lancé depuis le datacentre. »

Puis : « Je vous passe le relais dès que le paquet a traversé. »
""")

    # 9. Inventaire / uninstall / runtime
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "L'agent  ·  Points 4, 6, 10", "Il se paramètre, il se retire, il se montre", 9, "INTERNE")
    trips = [
        ("Plan de mesures", "Le centre pousse, dans la réponse au battement, ce que cet hôte doit mesurer. L'agent accuse. Plus de YAML édité à la main sur 128 postes."),
        ("Désinstallation", "Une commande locale. L'hôte prévient le centre, puis s'efface. Le parc ne garde pas des fantômes « en ligne »."),
        ("Où il tourne", "Service Windows, tâche, console. La fiche le dit. On ne cherche plus le processus à tâtons sur le poste."),
    ]
    for i, (title, body) in enumerate(trips):
        y = Inches(1.48) + i * Inches(1.75)
        rrect(s, MARGIN, y, Inches(12.2), Inches(1.6), WHITE, SLATE_200, adj=0.06)
        rrect(s, Inches(0.75), y + Inches(0.35), Inches(0.7), Inches(0.7), GOLD if i == 0 else (SKY if i == 1 else EMERALD), adj=0.18)
        txt(s, Inches(0.75), y + Inches(0.35), Inches(0.7), Inches(0.7), str(i + 1), 20, WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.7), y + Inches(0.28), Inches(10.5), Inches(0.4), title, 18, NAVY, bold=True)
        txt(s, Inches(1.7), y + Inches(0.75), Inches(10.5), Inches(0.65), body, 15, SLATE_600)
    notes(s, """
INTERNE — 2 min. Puis passation, phrase écrite :

« Le paquet est arrivé. La machine a parlé. Maintenant : qu'est-ce qu'on en fait, au centre ? »

Vous vous levez / vous avancez. Elle recule d'un pas. Visible.
""")

    # 10. Section plateforme
    s = section_slide(
        prs, "Relais 2", "Le centre — ce que l'opérateur voit",
        "Réception, vérité d'inventaire, seuils, courriel, qui a le droit de quoi.",
        10, "VOUS",
    )
    notes(s, "VOUS — 8 s. « Merci. On est de retour dans la salle. L'hôte n'est plus un fichier. C'est une fiche. »")

    # 11. Architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "La plateforme", "Ce qui tourne une fois le message reçu", 11, "VOUS")
    boxes = [
        ("API", "FastAPI\nauth agent + JWT humains"),
        ("Inventaire", "PostgreSQL\nfiches, alertes, audit"),
        ("Historique", "VictoriaMetrics\ncourbes dans le temps"),
        ("Console", "React\nsituation, parc, alertes"),
    ]
    for i, (title, body) in enumerate(boxes):
        x = MARGIN + i * Inches(3.15)
        fill = GOLD if i % 2 == 0 else NAVY
        fg = NAVY if fill == GOLD else WHITE
        rrect(s, x, Inches(1.55), Inches(2.95), Inches(2.15), fill, adj=0.08)
        txt(s, x + Inches(0.15), Inches(1.7), Inches(2.65), Inches(0.4), title, 16, fg, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.12), Inches(2.2), Inches(2.7), Inches(1.2), body, 14, fg, align=PP_ALIGN.CENTER)
    rrect(s, MARGIN, Inches(4.05), Inches(12.2), Inches(2.75), WHITE, SLATE_200, adj=0.05)
    txt(s, Inches(0.8), Inches(4.22), Inches(11.5), Inches(0.4), "Ce que ça change pour l'exploitation", 16, NAVY, bold=True)
    multiline(s, Inches(0.8), Inches(4.7), Inches(11.5), Inches(1.85), [
        "Les agents et les humains n'utilisent pas le même login. Clé d'agent ≠ mot de passe opérateur.",
        "Chaque écriture sensible laisse une trace (audit). On peut répondre « qui a révoqué cet hôte, et quand ».",
        "Docker Compose en laboratoire : Postgres, Redis, métriques, API, dashboard — une commande pour tout lever.",
        "On n'ouvre pas 6379/5432 sur le LAN. Boucle locale. La console parle à l'API, pas à la base.",
    ], size=15, color=SLATE_700)
    notes(s, """
VOUS — 2 min. Skip Loki/n8n sauf question.

« Si on vous demande Kubernetes : Lot suivant. Aujourd'hui on prouve le métier, pas l'orchestrateur. »
""")

    # 12. Console
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "La console  ·  Points 2, 3, 9", "Quatre gestes d'un matin d'opérateur", 12, "VOUS")
    screens = [
        ("Tableau de bord", "Parc vivant / silencieux, critiques ouvertes, santé du canal mail. On sait en 5 secondes si la nuit a été calme."),
        ("Parc", "Chaque hôte : OS, dernier contact, jauges. On ouvre une fiche, pas un RDP « pour voir »."),
        ("Fiche hôte", "Ce qui est constaté (nom, OS) ne se « corrige » pas à la main. Ce qui est attribué (responsable, équipe, VLAN d'exploitation) se corrige. L'écart VLAN constaté / attribué se voit."),
        ("Alertes", "Vérifier, valider, prendre en charge. Un tiroir, un verdict, un nom. Plus une boîte mail collective sans propriétaire."),
    ]
    for i, (fr, body) in enumerate(screens):
        y = Inches(1.42) + i * Inches(1.32)
        rrect(s, MARGIN, y, Inches(12.2), Inches(1.20), WHITE, SLATE_200, adj=0.08)
        rrect(s, Inches(0.75), y + Inches(0.28), Inches(2.55), Inches(0.64), NAVY, adj=0.12)
        txt(s, Inches(0.75), y + Inches(0.28), Inches(2.55), Inches(0.64), fr, 14, WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(3.55), y + Inches(0.22), Inches(8.9), Inches(0.80), body, 14, SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
    notes(s, """
VOUS — 2 min 30. Si démo live : ces quatre clics. Sinon restez ici.

Rôles en une phrase : Admin enrôle et paramètre. Opérateur triage. Consultation lit.
""")

    # 13. Alertes mail
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "La plateforme  ·  Point 8", "Une alerte qui trouve quelqu'un", 13, "VOUS")
    card(s, MARGIN, Inches(1.48), Inches(6.05), Inches(5.35), "Règles", [
        "CPU, RAM, disque : pas un pic d'une seconde. Une durée. Un spike n'ouvre pas une crise.",
        "",
        "Quatre gravités. Cycle ouvert → pris en charge → résolu.",
        "",
        "Fenêtres de maintenance distinctes des horaires de poste : un PC éteint à 20 h n'est pas un incident.",
        "",
        "Gabarit de courriel par type de vérification — pas un mail générique « alerte système ».",
    ], ROSE)
    card(s, Inches(6.9), Inches(1.48), Inches(5.85), Inches(5.35), "Canaux", [
        "Canal Lot 1 : API Mail CBC (clé, statut du canal visible dans la console).",
        "",
        "En plus : webhook HMAC signé — pour n8n ou un outil interne, sans mot de passe en clair dans l'URL.",
        "",
        "Ce qui n'est pas vendu aujourd'hui : ServiceNow, Teams, PagerDuty. On a le tuyau propre. Le métier choisira l'outil.",
    ], SKY)
    notes(s, """
VOUS — 2 min.

« Une alerte sans propriétaire est un bruit. Une alerte avec verdict et nom est un travail. »

Si Mail API n'est pas branchée en prod : le dire. Laboratoire vs production.
""")

    # 14. 10 points
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "Preuve", "Dix points livrés — pas un roadmap", 14, "ENSEMBLE")
    rows = [
        ("1", "Enrôlement", "INTERNE", "Jeton unique, fiche unique"),
        ("2", "Champs éditables", "VOUS", "L'UI ne propose que ce que l'API accepte"),
        ("3", "Responsable / équipe", "VOUS", "Plus de « non attribué » en lecture seule"),
        ("4", "Désinstallation", "INTERNE", "L'hôte prévient avant de partir"),
        ("5", "Reprise de contact", "INTERNE", "Un battement ramène l'hôte vivant"),
        ("6", "Plan de mesures", "VOUS", "Poussé au battement, accusé"),
        ("7", "Métriques", "INTERNE", "Collecte réelle, inventaire"),
        ("8", "Mail / webhook", "VOUS", "Gabarit + signature HMAC"),
        ("9", "Prise en charge", "VOUS", "Verdict, attribution, tiroir"),
        ("10", "Runtime OS", "INTERNE", "Où et comment l'agent tourne"),
    ]
    rrect(s, MARGIN, Inches(1.38), Inches(12.2), Inches(0.42), NAVY, adj=0.04)
    for i, h in enumerate(["#", "Point", "Voix", "Preuve en une ligne"]):
        xs = [0.7, 1.4, 5.3, 7.4]
        ws = [0.55, 3.7, 1.9, 5.0]
        txt(s, Inches(xs[i]), Inches(1.38), Inches(ws[i]), Inches(0.42), h, 11, GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    for r, row in enumerate(rows):
        y = Inches(1.82) + r * Inches(0.50)
        bg = WHITE if r % 2 == 0 else SLATE_100
        rect(s, MARGIN, y, Inches(12.2), Inches(0.50), bg)
        txt(s, Inches(0.7), y, Inches(0.55), Inches(0.50), row[0], 12, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.4), y, Inches(3.7), Inches(0.50), row[1], 13, SLATE_800, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(5.3), y, Inches(1.9), Inches(0.50), row[2], 11, GOLD_DK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(7.4), y, Inches(5.0), Inches(0.50), row[3], 13, SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
    notes(s, """
ENSEMBLE — 2 min. Alternez vraiment : elle lit 1, 4, 5, 7, 10. Vous lisez 2, 3, 6, 8, 9.

Rythme staccato. Pas de développement. La table EST le message : on a livré.
""")

    # 15. Démo
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "Preuve vivante", "Trois minutes, pas un tour du produit", 15, "ENSEMBLE")
    steps = [
        ("1", "Console", "VOUS", "Login admin. Tableau de bord. Un chiffre : hôtes vus, alertes ouvertes."),
        ("2", "Enrôler", "INTERNE", "Jeton labo. Agent de démo. Le poste apparaît. Pas de magie : un POST /enroll."),
        ("3", "Fiche", "VOUS", "Constats verrouillés, responsable, runtime. Ouvrir une alerte, la prendre en charge."),
        ("4", "Arrêt", "INTERNE", "Couper l'agent. Le bandeau hors ligne. Relancer : un battement suffit."),
    ]
    for i, (n, title, who, body) in enumerate(steps):
        y = Inches(1.48) + i * Inches(1.28)
        rrect(s, MARGIN, y, Inches(12.2), Inches(1.15), WHITE, SLATE_200, adj=0.06)
        rrect(s, Inches(0.75), y + Inches(0.28), Inches(0.55), Inches(0.55), GOLD, adj=0.18)
        txt(s, Inches(0.75), y + Inches(0.28), Inches(0.55), Inches(0.55), n, 16, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.5), y + Inches(0.18), Inches(2.2), Inches(0.38), title, 16, NAVY, bold=True)
        txt(s, Inches(3.8), y + Inches(0.20), Inches(1.6), Inches(0.32), who, 11, GOLD_DK, bold=True)
        txt(s, Inches(1.5), y + Inches(0.58), Inches(10.8), Inches(0.42), body, 14, SLATE_600)
    notes(s, """
ENSEMBLE — 3 min chrono. Si Docker est down : sautez, cette slide raconte déjà le film.

URL : http://localhost:3000  admin@cbc.cm
Agent : .\\scripts\\run-test-agent.ps1 -Token demo-token-123

Plan B : captures. Ne déboguez jamais en séance.
""")

    # 16. Pas encore
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "Honnêteté", "Ce qui n'est pas dans cette restitution", 16, "VOUS")
    card(s, MARGIN, Inches(1.48), Inches(6.05), Inches(5.35), "Pas encore — et on le dit", [
        "Rejeu de 24 h de métriques après une longue coupure (tampon disque propre).",
        "",
        "Actions distantes (redémarrer un service) : volontairement refusées. L'hôte n'est pas un shell.",
        "",
        "SSO entreprise, n8n en boucle fermée, ITSM, HA 5 000 hôtes.",
        "",
        "Listes officielles CBC (SWIFT, services métier) : placeholders, pas l'inventaire de production.",
    ], SLATE_700)
    card(s, Inches(6.9), Inches(1.48), Inches(5.85), Inches(5.35), "Ce dont on a besoin de vous", [
        "Fichier réseau sous-réseau → VLAN → libellé — pas un Excel « nom d'hôte / VLAN » qui périme au premier rebranchement.",
        "",
        "Listes de services et fichiers à surveiller, signées exploitation.",
        "",
        "URL et clé Mail CBC pour sortir du laboratoire.",
        "",
        "Un pilote : 5 à 10 machines, dont un serveur métier et un poste d'agence.",
    ], EMERALD)
    notes(s, """
VOUS — 2 min. C'est la slide qui vous gagne l'IT.

« On ne vous demande pas de croire. On vous demande un fichier, une liste, et dix machines. »
""")

    # 17. Close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, GOLD)
    txt(s, Inches(0.8), Inches(0.55), Inches(11.5), Inches(0.3), "À EMPORTER", 12, GOLD, bold=True)
    txt(s, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.7), "Quatre phrases", 32, WHITE, bold=True)
    takeaways = [
        "On remplace des scripts orphelins par un agent, un centre, une console.",
        "L'agent obtient un badge, mesure l'hôte, frappe en HTTPS — il n'ouvre aucun port.",
        "Le centre authentifie, historise, alerte, attribue. L'opérateur travaille sur des fiches, pas des RDP.",
        "Dix points sont livrés et démontrables. Le reste s'ajoute ; on ne le mime pas.",
    ]
    for i, line in enumerate(takeaways):
        y = Inches(1.85) + i * Inches(1.05)
        rrect(s, Inches(0.8), y, Inches(0.52), Inches(0.52), GOLD, adj=0.2)
        txt(s, Inches(0.8), y, Inches(0.52), Inches(0.52), str(i + 1), 16, NAVY, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, Inches(1.55), y, Inches(10.8), Inches(0.85), line, 16, WHITE, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.8), Inches(6.15), Inches(11.5), Inches(0.7),
        "Questions  ·  Agent → interne   ·   Console et règles → nous",
        16, SLATE_400)
    footer(s, 17, dark=True)
    notes(s, """
ENSEMBLE — 1 min. VOUS lisez 1 et 3. INTERNE lit 2. VOUS lit 4.

Silence après la quatrième. Puis : « Questions. »

Si plugin / IA / n8n : « Prochain chapitre. Aujourd'hui on a tenu le contrat : voir, collecter, alerter. »

Remerciez. Asseyez-vous en même temps.
""")

    # 18. Backup Q&A
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_chrome(s, "Si on vous coupe", "Réponses courtes — à deux", 18, "ENSEMBLE")
    qa = [
        ("C'est un EDR / un antivirus ?", "Non. On ne chasse pas le malware. On dit si la machine est vivante et chargée."),
        ("Vous ouvrez un port sur le poste ?", "Non. Sortant 443 uniquement. C'est un choix de sécurité, pas un raccourci."),
        ("On peut redémarrer un service à distance ?", "Pas maintenant. L'agent refuse. Quand ce sera le cas, il y aura une approbation humaine."),
        ("Et les scripts PowerShell ?", "On les éteint quand la même vérification est livrée et vue en production. Pas avant."),
        ("Qui a les droits ?", "Admin, opérateur, consultation. Les agents n'utilisent pas un compte utilisateur."),
        ("C'est en production CBC ?", "Laboratoire démontrable. Le pilote, c'est la prochaine décision — la vôtre."),
    ]
    for i, (q, a) in enumerate(qa):
        col, row = i % 2, i // 2
        x = MARGIN + col * Inches(6.25)
        y = Inches(1.42) + row * Inches(1.78)
        rrect(s, x, y, Inches(6.05), Inches(1.65), WHITE, SLATE_200, adj=0.06)
        txt(s, x + Inches(0.22), y + Inches(0.14), Inches(5.6), Inches(0.45), q, 14, NAVY, bold=True)
        txt(s, x + Inches(0.22), y + Inches(0.62), Inches(5.6), Inches(0.85), a, 13, SLATE_700)
    notes(s, """
Ne projetez cette slide que si le Q&A dérape. Sinon elle reste en réserve.

INTERNE prend ports / scripts / EDR. VOUS prend droits / prod / actions distantes.
""")

    out = Path(__file__).resolve().parent / "CBC_Supervision_Restitution_IT.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    print(build())
